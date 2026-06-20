"""
知识空间面板 — 创建/删除/上传文件/多选挂载
"""
import streamlit as st
from config import KB_CATEGORIES, KB_CATEGORY_LABELS


def render_kb_panel(kb_manager, active_kb_ids: list):
    """渲染知识空间管理面板"""
    st.subheader("📁 知识空间管理")

    # ── 创建新空间 ──
    with st.expander("➕ 创建新知识空间", expanded=False):
        name = st.text_input(
            "空间名称",
            placeholder="如：空气动力学、高数笔记...",
            key="kb_create_name",
            label_visibility="visible",
        )
        desc = st.text_area(
            "描述（可选）",
            placeholder="这个知识空间包含的内容...",
            key="kb_create_desc",
            label_visibility="visible",
        )
        # 改用 radio 代替 selectbox — 更容易看清
        st.markdown("**分类**")
        cat = st.radio(
            "分类",
            KB_CATEGORIES,
            format_func=lambda x: KB_CATEGORY_LABELS.get(x, x),
            key="kb_create_cat",
            label_visibility="collapsed",
            horizontal=True,
        )
        if st.button("✅ 创建知识空间", use_container_width=True, key="kb_create_btn"):
            if name.strip():
                kb_manager.create_kb(name.strip(), desc.strip(), cat)
                st.success(f"知识空间「{name}」创建成功！")
                st.rerun()
            else:
                st.warning("请输入空间名称")

    st.divider()

    # ── 知识空间列表 ──
    kbs = kb_manager.list_kbs()
    if not kbs:
        st.info("💡 还没有知识空间。点击上方「➕ 创建新知识空间」创建第一个！")
        st.caption("创建后可以上传 PDF 课件、TXT 笔记等，AI 将从这些资料中检索答案。")
        return active_kb_ids

    st.caption(f"共 **{len(kbs)}** 个知识空间 | 勾选空间后可上传文件")

    new_active = []
    for kb in kbs:
        cat_label = KB_CATEGORY_LABELS.get(kb["category"], "📋 其他")
        is_active = kb["id"] in active_kb_ids

        # ── 每个知识空间的卡片 ──
        with st.container():
            st.markdown("---") if kb != kbs[0] else None

            # 第一行：勾选框 + 名称 + 删除
            c1, c2 = st.columns([0.7, 0.3])
            with c1:
                checked = st.checkbox(
                    f"**{kb['name']}**  ·  {cat_label}  ·  {kb['doc_count']} 个文档",
                    value=is_active,
                    key=f"kb_select_{kb['id']}",
                )
                if checked:
                    new_active.append(kb["id"])
                st.caption(f"创建于 {kb.get('created_at', '')}  |  {kb.get('description', '')}")
            with c2:
                if st.button("🗑️ 删除此空间", key=f"kb_del_{kb['id']}", use_container_width=True):
                    kb_manager.delete_kb(kb["id"])
                    if kb["id"] in active_kb_ids:
                        active_kb_ids.remove(kb["id"])
                    st.rerun()

            # 第二行：文件管理（始终显示，不只是 checked 后）
            with st.expander(f"📂 管理「{kb['name']}」的文件 — 上传 PDF / TXT", expanded=checked and kb["doc_count"] == 0):
                _render_file_manager(kb_manager, kb["id"])

    return new_active


def _render_file_manager(kb_manager, kb_id: str):
    """文件上传与删除"""
    kb_info = kb_manager.get_kb_info(kb_id)
    if not kb_info:
        return

    # 上传区
    st.markdown("**📤 上传新文件**")
    uploaded = st.file_uploader(
        "选择 PDF、PPT、TXT 或 Markdown 文件",
        type=["pdf", "pptx", "ppt", "txt", "md"],
        key=f"file_upload_{kb_id}",
        label_visibility="visible",
    )
    if uploaded:
        with st.spinner(f"正在解析「{uploaded.name}」..."):
            try:
                text = _parse_uploaded_file(uploaded)
                if text.strip():
                    count = kb_manager.add_file(kb_id, uploaded.name, text.strip())
                    st.success(f"✅ 已入库「{uploaded.name}」→ {count} 个文本块")
                    st.rerun()
                else:
                    st.warning("⚠️ 文件中未提取到有效文字。如果是扫描版 PDF，请先 OCR 识别。")
            except Exception as e:
                st.error(f"解析失败：{str(e)[:200]}")

    # 已有文件列表
    st.divider()
    files = kb_manager.list_files(kb_id)
    if files:
        st.markdown(f"**📋 已上传 {len(files)} 个文件**")
        for f in files:
            fc1, fc2 = st.columns([0.85, 0.15])
            with fc1:
                st.markdown(
                    f"📄 **{f['file_name']}** · {f['chunk_count']} 块 · {f['uploaded_at']}"
                )
            with fc2:
                if st.button("🗑️", key=f"file_del_{kb_id}_{f['file_name']}", help=f"删除 {f['file_name']}"):
                    kb_manager.delete_file(kb_id, f["file_name"])
                    st.rerun()
    else:
        st.info("📭 此空间暂无文件。上传一个 PDF 课件或 TXT 笔记开始使用。")


def _parse_uploaded_file(uploaded) -> str:
    """
    解析上传的文件，返回提取的纯文本。

    支持格式：PDF、PPTX/PPT、TXT、Markdown
    """
    fname = uploaded.name.lower()

    if fname.endswith(".pdf"):
        import fitz
        doc = fitz.open(stream=uploaded.read(), filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text

    if fname.endswith((".pptx", ".ppt")):
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
        try:
            prs = Presentation(uploaded)
        except PackageNotFoundError:
            # 旧版 .ppt 格式无法用 python-pptx 解析
            return ""
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[幻灯片 {slide_num}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_parts.append(row_text)
            slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)

    # TXT / MD / 其他纯文本格式
    return uploaded.read().decode("utf-8", errors="ignore")
